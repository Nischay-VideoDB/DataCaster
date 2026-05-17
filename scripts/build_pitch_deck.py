"""Build the DataCaster pitch deck (4 slides).

Outputs: ./DataCaster_Pitch.pptx at the repo root.

Slide flow (sets the stage before the live demo):
  1. Title — hero line, almost no body copy
  2. Problem — manual scouting today: cause → effect
  3. Solution — DataCaster: one feed in, three data products out
  4. Why it wins — rubric mapping + CTA

Uses the existing VideoDB wordmark from frontend/public/videodb-logo.png and
synthesises a DataCaster mark from PIL primitives so we don't need to fetch
any external assets.
"""

from __future__ import annotations

import math
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Palette — sampled from the DataCaster favicon + VideoDB wordmark.

ZINC_950 = RGBColor(0x09, 0x09, 0x0B)
ZINC_900 = RGBColor(0x18, 0x18, 0x1B)
ZINC_850 = RGBColor(0x1F, 0x1F, 0x23)
ZINC_800 = RGBColor(0x27, 0x27, 0x2A)
ZINC_700 = RGBColor(0x3F, 0x3F, 0x46)
ZINC_500 = RGBColor(0x71, 0x71, 0x7A)
ZINC_400 = RGBColor(0xA1, 0xA1, 0xAA)
ZINC_300 = RGBColor(0xD4, 0xD4, 0xD8)
ZINC_200 = RGBColor(0xE4, 0xE4, 0xE7)
ZINC_100 = RGBColor(0xF4, 0xF4, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xEC, 0x5B, 0x16)
ORANGE_SOFT = RGBColor(0xF9, 0x73, 0x16)
RED = RGBColor(0xEF, 0x44, 0x44)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
SKY = RGBColor(0x38, 0xBD, 0xF8)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)

REPO = Path(__file__).resolve().parent.parent
ASSETS = REPO / "scripts" / "_pitch_assets"
ASSETS.mkdir(parents=True, exist_ok=True)
OUT = REPO / "DataCaster_Pitch.pptx"

VIDEODB_LOGO = REPO / "frontend" / "public" / "videodb-logo.png"


# ---------------------------------------------------------------------------
# Asset synthesis

def _font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/SFNSDisplay.ttf",
        "/System/Library/Fonts/HelveticaNeue.ttc",
    ]
    for c in candidates:
        if Path(c).exists():
            try:
                return ImageFont.truetype(c, size)
            except OSError:
                continue
    return ImageFont.load_default()


def render_datacaster_mark(path: Path, size: int = 512) -> Path:
    """Square dark-tile + orange radio glyph, matching favicon.svg."""
    img = Image.new("RGBA", (size, size), (9, 9, 11, 255))
    radius = int(size * 0.18)
    mask = Image.new("L", (size, size), 0)
    ImageDraw.Draw(mask).rounded_rectangle((0, 0, size - 1, size - 1), radius=radius, fill=255)
    img.putalpha(mask)
    draw = ImageDraw.Draw(img)

    cx = cy = size // 2
    stroke = max(2, size // 40)
    color = (0xEC, 0x5B, 0x16, 255)

    def arc(radius_px: int, span_deg: int) -> None:
        bbox = (cx - radius_px, cy - radius_px, cx + radius_px, cy + radius_px)
        draw.arc(bbox, start=-span_deg, end=span_deg, fill=color, width=stroke)
        draw.arc(bbox, start=180 - span_deg, end=180 + span_deg, fill=color, width=stroke)

    arc(int(size * 0.20), 50)
    arc(int(size * 0.32), 55)

    dot_r = max(6, size // 28)
    draw.ellipse((cx - dot_r, cy - dot_r, cx + dot_r, cy + dot_r), fill=color)

    img.save(path, "PNG")
    return path


# ---------------------------------------------------------------------------
# Slide construction helpers.

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _emu(x: float) -> Emu:
    return Inches(x)


def add_bg(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False


def add_text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: float = 18,
    bold: bool = False,
    color: RGBColor = WHITE,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font_name: str = "Helvetica Neue",
    line_spacing: float | None = None,
):
    tb = slide.shapes.add_textbox(_emu(left), _emu(top), _emu(width), _emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.0)
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: RGBColor = ZINC_900,
    border: RGBColor | None = ZINC_800,
    corner: float = 0.18,
    border_width: float = 0.75,
):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, _emu(left), _emu(top), _emu(width), _emu(height)
    )
    card.adjustments[0] = corner
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if border is None:
        card.line.fill.background()
    else:
        card.line.color.rgb = border
        card.line.width = Pt(border_width)
    return card


def add_image(slide, path: Path, left: float, top: float, width: float | None = None, height: float | None = None):
    kwargs = {}
    if width is not None:
        kwargs["width"] = _emu(width)
    if height is not None:
        kwargs["height"] = _emu(height)
    return slide.shapes.add_picture(str(path), _emu(left), _emu(top), **kwargs)


def add_brand_header(slide, dc_mark: Path) -> None:
    """Top-left DataCaster brand bug."""
    add_image(slide, dc_mark, left=0.55, top=0.4, height=0.55)
    add_text(
        slide, left=1.2, top=0.45, width=4.0, height=0.5,
        text="DataCaster", size=18, bold=True, color=ZINC_100,
    )


def add_brand_footer(slide) -> None:
    """Bottom-right 'powered by VideoDB'."""
    add_text(
        slide, left=8.5, top=7.05, width=2.5, height=0.35,
        text="powered by",
        size=9, color=ZINC_500, align=PP_ALIGN.RIGHT, bold=True,
    )
    add_image(slide, VIDEODB_LOGO, left=11.1, top=6.98, height=0.32)


# ---------------------------------------------------------------------------
# Slide 1 — Title (lean).

def slide_title(prs: Presentation, dc_mark: Path) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, ZINC_950)

    add_image(s, dc_mark, left=2.7, top=2.3, height=1.4)
    add_text(
        s, left=4.3, top=2.3, width=7.0, height=1.4,
        text="DataCaster",
        size=72, bold=True, color=ZINC_100,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    bar = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, _emu(6.27), _emu(3.95), _emu(0.8), _emu(0.07)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE

    add_text(
        s, left=1.0, top=4.2, width=11.3, height=0.85,
        text="One stream in, three sportsbook-grade products out.",
        size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )

    # Powered-by row at the bottom.
    add_text(
        s, left=4.5, top=6.85, width=2.5, height=0.35,
        text="powered by",
        size=11, color=ZINC_500, align=PP_ALIGN.RIGHT, bold=True,
    )
    add_image(s, VIDEODB_LOGO, left=7.15, top=6.78, height=0.42)


# ---------------------------------------------------------------------------
# Slide 2 — Problem.

def slide_problem(prs: Presentation, dc_mark: Path) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, ZINC_950)
    add_brand_header(s, dc_mark)

    add_text(
        s, left=0.55, top=1.15, width=12.5, height=0.45,
        text="THE PROBLEM",
        size=12, bold=True, color=ORANGE,
    )
    add_text(
        s, left=0.55, top=1.55, width=12.5, height=0.7,
        text="Live sports data is still typed by humans.",
        size=32, bold=True, color=WHITE,
    )
    add_text(
        s, left=0.55, top=2.3, width=12.5, height=0.5,
        text="Every goal, card, and corner you see in a betting app started with a person watching a screen.",
        size=14, color=ZINC_400,
    )

    # Three vertically aligned cards: cause → effect → cost.
    cards = [
        {
            "tag": "CAUSE",
            "color": RED,
            "title": "Manual scouting rooms",
            "body": (
                "Sportradar and Stats Perform pay rooms full\n"
                "of analysts to type events into databases\n"
                "in real time, frame by frame."
            ),
        },
        {
            "tag": "EFFECT",
            "color": AMBER,
            "title": "Slow, expensive, error-prone",
            "body": (
                "Latency in seconds. Coverage gaps for tier-3\n"
                "leagues. Data mismatches across vendors.\n"
                "Highlights still cut hours after the whistle."
            ),
        },
        {
            "tag": "COST",
            "color": ORANGE,
            "title": "Most matches stay dark",
            "body": (
                "If a match isn't worth a scout, it gets no live\n"
                "data, no search, no clips — invisible to\n"
                "sportsbooks, broadcasters, and fans."
            ),
        },
    ]

    card_w = 4.05
    card_h = 3.45
    gap = 0.18
    start_left = (13.333 - (card_w * 3 + gap * 2)) / 2
    top = 3.05

    for i, c in enumerate(cards):
        left = start_left + i * (card_w + gap)
        add_card(s, left, top, card_w, card_h, fill=ZINC_900, border=ZINC_800)

        bar = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, _emu(left), _emu(top + 0.4),
            _emu(0.08), _emu(card_h - 0.8),
        )
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = c["color"]

        chip_w = 0.95
        chip = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            _emu(left + 0.4), _emu(top + 0.45),
            _emu(chip_w), _emu(0.4),
        )
        chip.adjustments[0] = 0.5
        chip.fill.solid()
        chip.fill.fore_color.rgb = ZINC_850
        chip.line.color.rgb = c["color"]
        chip.line.width = Pt(0.75)
        add_text(
            s, left=left + 0.4, top=top + 0.45, width=chip_w, height=0.4,
            text=c["tag"], size=10, bold=True, color=c["color"],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        add_text(
            s, left=left + 0.4, top=top + 1.1, width=card_w - 0.6, height=0.7,
            text=c["title"], size=18, bold=True, color=WHITE,
        )
        add_text(
            s, left=left + 0.4, top=top + 1.95, width=card_w - 0.6, height=1.4,
            text=c["body"], size=12.5, color=ZINC_300, line_spacing=1.4,
        )

    add_text(
        s, left=0.55, top=6.7, width=12.2, height=0.45,
        text="Sports broadcast was never really live data.  ↓  We made it live.",
        size=14, bold=True, color=ZINC_300, align=PP_ALIGN.CENTER,
    )

    add_brand_footer(s)


# ---------------------------------------------------------------------------
# Slide 3 — Solution (the architecture).

def slide_solution(prs: Presentation, dc_mark: Path) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, ZINC_950)
    add_brand_header(s, dc_mark)

    add_text(
        s, left=0.55, top=1.15, width=12.5, height=0.45,
        text="THE SOLUTION",
        size=12, bold=True, color=ORANGE,
    )
    add_text(
        s, left=0.55, top=1.55, width=12.5, height=0.7,
        text="One video feed → three production data products.",
        size=28, bold=True, color=WHITE,
    )

    # Pipeline strip: input pill → engine → outputs.
    pill = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, _emu(0.55), _emu(2.55), _emu(2.6), _emu(0.7)
    )
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = ZINC_800
    pill.line.color.rgb = ZINC_700
    add_text(
        s, left=0.55, top=2.55, width=2.6, height=0.7,
        text="Video feed in",
        size=14, bold=True, color=ZINC_100,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        s, left=0.55, top=3.3, width=2.6, height=0.4,
        text="RTSP · RTMP · YouTube",
        size=11, color=ZINC_400, align=PP_ALIGN.CENTER,
    )

    arrow = s.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, _emu(3.25), _emu(2.75), _emu(0.7), _emu(0.3)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ORANGE
    arrow.line.fill.background()

    engine = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, _emu(4.0), _emu(2.5), _emu(4.6), _emu(0.85)
    )
    engine.adjustments[0] = 0.3
    engine.fill.solid()
    engine.fill.fore_color.rgb = ZINC_900
    engine.line.color.rgb = ORANGE
    engine.line.width = Pt(1.5)
    add_text(
        s, left=4.0, top=2.5, width=4.6, height=0.5,
        text="VideoDB pipeline", size=16, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        s, left=4.0, top=2.9, width=4.6, height=0.45,
        text="RTStream · index_visuals · index_audio · transcripts",
        size=11, color=ORANGE_SOFT, align=PP_ALIGN.CENTER,
    )

    arrow2 = s.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, _emu(8.7), _emu(2.75), _emu(0.7), _emu(0.3)
    )
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = ORANGE
    arrow2.line.fill.background()

    add_text(
        s, left=9.5, top=2.55, width=3.8, height=0.7,
        text="3 data products", size=18, bold=True, color=ORANGE,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    cards = [
        {
            "title": "Structured event JSON",
            "accent": EMERALD,
            "icon": "▣",
            "lines": [
                "Goals · saves · cards · corners",
                "Confidence + team labels",
                "SSE feed + per-video SQLite",
            ],
            "vdb": "video.index_scenes · search",
        },
        {
            "title": "Searchable memory",
            "accent": SKY,
            "icon": "🔍",
            "lines": [
                "Natural-language Q&A",
                "Visual + audio + transcript",
                "Query expansion bypasses floor",
            ],
            "vdb": "video.search · generate_text",
        },
        {
            "title": "Programmable editing",
            "accent": VIOLET,
            "icon": "✂",
            "lines": [
                "9:16 highlight reel · 1-click",
                "30s recap caption · OmniVoice",
                "Auto-posted to Telegram",
            ],
            "vdb": "Timeline · generate_voice",
        },
    ]
    card_w = 4.05
    card_h = 3.05
    gap = 0.18
    start_left = (13.333 - (card_w * 3 + gap * 2)) / 2
    top = 3.85

    for i, c in enumerate(cards):
        left = start_left + i * (card_w + gap)
        add_card(s, left, top, card_w, card_h, fill=ZINC_900, border=ZINC_800)
        stripe = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, _emu(left + 0.25), _emu(top + 0.3), _emu(0.45), _emu(0.08)
        )
        stripe.line.fill.background()
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = c["accent"]
        add_text(
            s, left=left + 0.25, top=top + 0.42, width=card_w - 0.5, height=0.5,
            text=f"{c['icon']}  {c['title']}", size=18, bold=True, color=WHITE,
        )
        body = "\n".join("•  " + line for line in c["lines"])
        add_text(
            s, left=left + 0.25, top=top + 1.05, width=card_w - 0.5, height=1.4,
            text=body, size=12, color=ZINC_300, line_spacing=1.4,
        )
        add_text(
            s, left=left + 0.25, top=top + 2.45, width=card_w - 0.5, height=0.45,
            text=c["vdb"], size=10, bold=True, color=c["accent"],
        )

    add_text(
        s, left=0.55, top=7.05, width=8.0, height=0.35,
        text="Re-runs are free  ·  events persist by video_id  ·  every component degrades gracefully",
        size=10, color=ZINC_500,
    )
    add_brand_footer(s)


# ---------------------------------------------------------------------------
# Slide 4 — Why it wins (rubric mapping + CTA).

def slide_why(prs: Presentation, dc_mark: Path) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, ZINC_950)
    add_brand_header(s, dc_mark)

    add_text(
        s, left=0.55, top=1.15, width=12.5, height=0.45,
        text="WHY IT WINS",
        size=12, bold=True, color=ORANGE,
    )
    add_text(
        s, left=0.55, top=1.55, width=12.5, height=0.7,
        text="Built for the rubric, not for the demo.",
        size=28, bold=True, color=WHITE,
    )

    rubric = [
        {
            "weight": "40%",
            "title": "Technical execution",
            "color": EMERALD,
            "points": [
                "Every component independently demoable",
                "Graceful degradation: search works if commentary fails",
                "Test runner covers VOD + describe + live RTStream",
            ],
        },
        {
            "weight": "30%",
            "title": "Creativity",
            "color": ORANGE,
            "points": [
                "Automates Sportradar's manual scouting loop",
                "One pipeline → three data products",
                "Two classifier modes: football + describe-scenes",
            ],
        },
        {
            "weight": "30%",
            "title": "Depth of VideoDB",
            "color": SKY,
            "points": [
                "RTStream + index_visuals + index_audio + transcripts",
                "Search · generate_text · OmniVoice · Timeline",
                "Sandbox lifecycle managed in finally:",
            ],
        },
    ]

    card_w = 4.05
    card_h = 3.55
    gap = 0.18
    start_left = (13.333 - (card_w * 3 + gap * 2)) / 2
    top = 2.45

    for i, r in enumerate(rubric):
        left = start_left + i * (card_w + gap)
        add_card(s, left, top, card_w, card_h, fill=ZINC_900, border=ZINC_800)

        badge_w = 1.0
        wb = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            _emu(left + card_w - badge_w - 0.25), _emu(top + 0.3),
            _emu(badge_w), _emu(0.45),
        )
        wb.adjustments[0] = 0.5
        wb.fill.solid()
        wb.fill.fore_color.rgb = ZINC_800
        wb.line.color.rgb = r["color"]
        add_text(
            s, left=left + card_w - badge_w - 0.25, top=top + 0.3,
            width=badge_w, height=0.45,
            text=r["weight"], size=13, bold=True, color=r["color"],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        stripe = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, _emu(left + 0.3), _emu(top + 0.42), _emu(0.45), _emu(0.08)
        )
        stripe.line.fill.background()
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = r["color"]

        add_text(
            s, left=left + 0.3, top=top + 0.55, width=card_w - 1.4, height=0.6,
            text=r["title"], size=20, bold=True, color=WHITE,
        )

        body = "\n".join("•  " + line for line in r["points"])
        add_text(
            s, left=left + 0.3, top=top + 1.4, width=card_w - 0.6, height=2.0,
            text=body, size=12.5, color=ZINC_300, line_spacing=1.45,
        )

    cta = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        _emu(0.55), _emu(6.4), _emu(12.22), _emu(0.8),
    )
    cta.adjustments[0] = 0.4
    cta.fill.solid()
    cta.fill.fore_color.rgb = ZINC_900
    cta.line.color.rgb = ORANGE
    cta.line.width = Pt(1.5)
    add_text(
        s, left=0.55, top=6.4, width=12.22, height=0.8,
        text="▶  Now let's see it run.",
        size=22, bold=True, color=ORANGE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )


# ---------------------------------------------------------------------------
# Main.

def main() -> None:
    dc_mark = render_datacaster_mark(ASSETS / "datacaster_mark.png", size=512)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, dc_mark)
    slide_problem(prs, dc_mark)
    slide_solution(prs, dc_mark)
    slide_why(prs, dc_mark)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
